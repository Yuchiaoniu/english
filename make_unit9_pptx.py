# -*- coding: utf-8 -*-
"""Unit 9 — How Do You Feel? 單課完整簡報產生器。
9 個情緒/感受單字 + 連連看 + 句型填空 + 會話，共 15 頁。
輸出：Unit9_How_Do_You_Feel.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------- 主題色（rose 玫瑰紅，搭情緒主題）----------
THEME   = RGBColor(0xE0, 0x4A, 0x6B)
THEME_D = RGBColor(0xB0, 0x35, 0x52)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
DARK    = RGBColor(0x2C, 0x3E, 0x50)
GRAY    = RGBColor(0x7F, 0x8C, 0x8D)
LIGHT   = RGBColor(0xFB, 0xF4, 0xF6)
A_BLUE  = RGBColor(0x4A, 0x90, 0xD9)
B_ORNG  = RGBColor(0xF0, 0x8A, 0x24)
GREEN   = RGBColor(0x27, 0xAE, 0x60)
CARD_BD = RGBColor(0xEA, 0xD4, 0xDB)

W = Inches(13.333)
H = Inches(7.5)

UNIT_NO    = 9
UNIT_EN    = "How Do You Feel?"
UNIT_ZH    = "你感覺怎麼樣？"
UNIT_EMOJI = "💖"

VOCAB = [
    {"word": "happy",     "zh": "快樂",   "zhuyin": "ㄎㄨㄞˋ ㄌㄜˋ",      "emoji": "😄"},
    {"word": "sad",       "zh": "難過",   "zhuyin": "ㄋㄢˊ ㄍㄨㄛˋ",      "emoji": "😢"},
    {"word": "angry",     "zh": "生氣",   "zhuyin": "ㄕㄥ ㄑㄧˋ",         "emoji": "😠"},
    {"word": "scared",    "zh": "害怕",   "zhuyin": "ㄏㄞˋ ㄆㄚˋ",        "emoji": "😨"},
    {"word": "shocked",   "zh": "震驚",   "zhuyin": "ㄓㄣˋ ㄐㄧㄥ",        "emoji": "😲"},
    {"word": "sick",      "zh": "生病",   "zhuyin": "ㄕㄥ ㄅㄧㄥˋ",       "emoji": "🤒"},
    {"word": "proud",     "zh": "驕傲",   "zhuyin": "ㄐㄧㄠ ㄠˋ",         "emoji": "🏆"},
    {"word": "confident", "zh": "有自信", "zhuyin": "ㄧㄡˇ ㄗˋ ㄒㄧㄣˋ",  "emoji": "💪"},
    {"word": "cool",      "zh": "酷",     "zhuyin": "ㄎㄨˋ",              "emoji": "😎"},
]

# ---------- 連連看一：單字 ↔ 中文 ----------
MATCH1_LEFT = [
    ("A", "scared"),
    ("B", "happy"),
    ("C", "sick"),
    ("D", "angry"),
    ("E", "sad"),
]
MATCH1_RIGHT = [
    ("1", "難過"),
    ("2", "害怕"),
    ("3", "生氣"),
    ("4", "生病"),
    ("5", "快樂"),
]
MATCH1_ANS = "A–2  ·  B–5  ·  C–4  ·  D–3  ·  E–1"

# ---------- 連連看二：情境 ↔ 情緒詞 ----------
MATCH2_LEFT = [
    ("A", "I got 100 in the test!"),
    ("B", "Surprise! Happy birthday!"),
    ("C", "I'll win the race!"),
    ("D", "Look at my sunglasses!"),
]
MATCH2_RIGHT = [
    ("1", "cool"),
    ("2", "confident"),
    ("3", "proud"),
    ("4", "shocked"),
]
MATCH2_ANS = "A–3  ·  B–4  ·  C–2  ·  D–1"

# ---------- 句型 + 填空 ----------
PATTERN1 = {
    "structure": "I feel ___.       /       I am ___.",
    "zh": "我覺得 ___ 。  /  我 ___ 。",
    "examples": [
        ("I feel happy when I play games.",     "我玩遊戲時覺得很快樂。"),
        ("I am sad because I lost my book.",    "我很難過因為我弄丟了書。"),
    ],
    "fill": [
        {"q": "I feel ___ when my team wins.",
         "opts": ["proud", "sad", "sick"],   "ans": 0},
        {"q": "I am ___ because the dog is so big!",
         "opts": ["happy", "scared", "cool"], "ans": 1},
    ],
}

PATTERN2 = {
    "structure": "He / She looks ___.       /       Why are you ___?",
    "zh": "他／她看起來 ___ 。  /  你為什麼 ___ ？",
    "examples": [
        ("She looks angry today.",  "她今天看起來在生氣。"),
        ("Why are you so happy?",   "你為什麼這麼快樂？"),
    ],
    "fill": [
        {"q": "He looks ___. He has a fever.",
         "opts": ["cool", "sick", "proud"],   "ans": 1},
        {"q": "Why are you so ___? You won the contest!",
         "opts": ["sad", "proud", "angry"],   "ans": 1},
    ],
}

# ---------- 會話 ----------
DIALOGUE = [
    ("A", "Sue", "Hi, Ken! How are you today?",                "嗨 Ken！你今天怎麼樣？"),
    ("B", "Ken", "I feel a little sick.",                      "我覺得有點不舒服。"),
    ("A", "Sue", "Oh no! I'm sad to hear that.",               "喔不！聽到這個我很難過。"),
    ("B", "Ken", "But guess what? I got 100 in math!",         "但你猜怎麼著？我數學考了 100 分！"),
    ("A", "Sue", "Wow! That's so cool! You should be proud!",  "哇！超酷的！你應該很驕傲！"),
    ("B", "Ken", "I am! And I feel confident and happy now.",  "對啊！而且現在我很有自信也很快樂。"),
]


# ---------- 共用工具 ----------
def slide_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def textbox(slide, x, y, w, h, lines, anchor=MSO_ANCHOR.MIDDLE, wrap=True):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for i, (text, size, bold, color, align) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return box


def rrect(slide, x, y, w, h, fill, line=None, line_w=1.0):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.shadow.inherit = False
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
    return sh


def topbar(slide, label):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.62))
    bar.shadow.inherit = False
    bar.fill.solid()
    bar.fill.fore_color.rgb = THEME
    bar.line.fill.background()
    tf = bar.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "   Unit %d · %s   %s   %s" % (UNIT_NO, UNIT_EN, UNIT_EMOJI, label)
    r.font.size = Pt(17)
    r.font.bold = True
    r.font.color.rgb = WHITE


def heading(slide, text):
    textbox(slide, Inches(0.5), Inches(0.72), Inches(12.3), Inches(0.85),
            [(text, 32, True, THEME_D, PP_ALIGN.LEFT)])


# ---------- 投影片 ----------
def s_cover(prs):
    s = slide_blank(prs)
    set_bg(s, THEME)
    deco = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.3), Inches(-1.4),
                              Inches(4.6), Inches(4.6))
    deco.shadow.inherit = False
    deco.fill.solid(); deco.fill.fore_color.rgb = THEME_D
    deco.line.fill.background()
    deco2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.6), Inches(4.6),
                               Inches(4.2), Inches(4.2))
    deco2.shadow.inherit = False
    deco2.fill.solid(); deco2.fill.fore_color.rgb = THEME_D
    deco2.line.fill.background()

    textbox(s, Inches(0.8), Inches(1.4), Inches(11.7), Inches(1.7),
            [(UNIT_EMOJI, 96, False, WHITE, PP_ALIGN.CENTER)])
    textbox(s, Inches(0.8), Inches(3.05), Inches(11.7), Inches(1.2),
            [("Unit %d · %s" % (UNIT_NO, UNIT_EN), 48, True, WHITE, PP_ALIGN.CENTER)])
    textbox(s, Inches(0.8), Inches(4.25), Inches(11.7), Inches(0.9),
            [(UNIT_ZH, 40, True, WHITE, PP_ALIGN.CENTER)])
    textbox(s, Inches(0.8), Inches(5.45), Inches(11.7), Inches(0.7),
            [("四年級英語 · 單字 9 · 連連看 · 句型填空 · 會話", 22, False, WHITE, PP_ALIGN.CENTER)])
    textbox(s, Inches(0.8), Inches(6.35), Inches(11.7), Inches(0.55),
            [("🌐  https://Yuchiaoniu.github.io/english/", 18, False, WHITE, PP_ALIGN.CENTER)])


def s_vocab_overview(prs):
    s = slide_blank(prs)
    set_bg(s, LIGHT)
    topbar(s, "Vocabulary 單字總覽")
    heading(s, "本課單字 9 個")
    cols, rows = 3, 3
    cw, ch = Inches(3.95), Inches(1.55)
    gx, gy = Inches(0.2), Inches(0.18)
    sx, sy = Inches(0.46), Inches(1.75)
    for i, v in enumerate(VOCAB):
        c = i % cols
        r = i // cols
        x = sx + c * (cw + gx)
        y = sy + r * (ch + gy)
        card = rrect(s, x, y, cw, ch, WHITE, line=CARD_BD, line_w=1.5)
        tf = card.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        p1 = tf.paragraphs[0]; p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run(); r1.text = "%s  %s" % (v["emoji"], v["word"])
        r1.font.size = Pt(24); r1.font.bold = True; r1.font.color.rgb = DARK
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = v["zh"]
        r2.font.size = Pt(20); r2.font.color.rgb = THEME


def s_vocab_cards(prs, items, label, title):
    s = slide_blank(prs)
    set_bg(s, LIGHT)
    topbar(s, label)
    heading(s, title)
    n = len(items)
    cw_in = 4.0
    ch_in = 4.3
    gap = 0.3
    total_w = n * cw_in + (n - 1) * gap
    sx_in = (13.333 - total_w) / 2
    sy = Inches(1.95)
    for i, v in enumerate(items):
        x = Inches(sx_in + i * (cw_in + gap))
        card = rrect(s, x, sy, Inches(cw_in), Inches(ch_in),
                     WHITE, line=THEME, line_w=1.75)
        tf = card.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        p0 = tf.paragraphs[0]; p0.alignment = PP_ALIGN.CENTER
        r0 = p0.add_run(); r0.text = v["emoji"]; r0.font.size = Pt(72)
        p1 = tf.add_paragraph(); p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run(); r1.text = v["word"]
        r1.font.size = Pt(32); r1.font.bold = True; r1.font.color.rgb = THEME_D
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = v["zh"]
        r2.font.size = Pt(26); r2.font.bold = True; r2.font.color.rgb = DARK
        p3 = tf.add_paragraph(); p3.alignment = PP_ALIGN.CENTER
        r3 = p3.add_run(); r3.text = v["zhuyin"]
        r3.font.size = Pt(18); r3.font.color.rgb = GRAY


def s_section(prs, emoji, en, zh):
    s = slide_blank(prs)
    set_bg(s, THEME)
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.55), W, Inches(2.4))
    band.shadow.inherit = False
    band.fill.solid(); band.fill.fore_color.rgb = THEME_D
    band.line.fill.background()
    textbox(s, Inches(0.8), Inches(1.15), Inches(11.7), Inches(1.4),
            [(emoji, 80, False, WHITE, PP_ALIGN.CENTER)])
    textbox(s, Inches(0.8), Inches(2.75), Inches(11.7), Inches(1.1),
            [(en, 48, True, WHITE, PP_ALIGN.CENTER)])
    textbox(s, Inches(0.8), Inches(3.9), Inches(11.7), Inches(0.9),
            [(zh, 32, False, WHITE, PP_ALIGN.CENTER)])


def s_match(prs, label, title, instr, left, right):
    s = slide_blank(prs)
    set_bg(s, LIGHT)
    topbar(s, label)
    heading(s, title)
    textbox(s, Inches(0.5), Inches(1.55), Inches(12.3), Inches(0.5),
            [(instr, 18, False, GRAY, PP_ALIGN.LEFT)])
    n = max(len(left), len(right))
    if n == 5:
        row_h, gap, start_y = 0.85, 0.12, 2.15
    else:
        row_h, gap, start_y = 1.05, 0.18, 2.35
    col_w = Inches(4.8)
    left_x = Inches(1.0)
    right_x = Inches(7.55)
    for i, (lbl, txt) in enumerate(left):
        y = start_y + i * (row_h + gap)
        card = rrect(s, left_x, Inches(y), col_w, Inches(row_h),
                     WHITE, line=THEME, line_w=1.5)
        tf = card.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r1 = p.add_run(); r1.text = "  " + lbl + ".   "
        r1.font.size = Pt(22); r1.font.bold = True; r1.font.color.rgb = THEME
        r2 = p.add_run(); r2.text = txt
        r2.font.size = Pt(22); r2.font.bold = True; r2.font.color.rgb = DARK
    for i, (lbl, txt) in enumerate(right):
        y = start_y + i * (row_h + gap)
        card = rrect(s, right_x, Inches(y), col_w, Inches(row_h),
                     WHITE, line=THEME, line_w=1.5)
        tf = card.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r1 = p.add_run(); r1.text = "  " + lbl + ".   "
        r1.font.size = Pt(22); r1.font.bold = True; r1.font.color.rgb = THEME
        r2 = p.add_run(); r2.text = txt
        r2.font.size = Pt(22); r2.font.bold = True; r2.font.color.rgb = DARK


def s_pattern_fill(prs, idx, p):
    s = slide_blank(prs)
    set_bg(s, LIGHT)
    topbar(s, "Pattern %d/2 句型填空" % idx)
    heading(s, "句型 %d" % idx)
    # 結構列
    box1 = rrect(s, Inches(0.6), Inches(1.62), Inches(12.13), Inches(0.95),
                 THEME, line=None)
    tf1 = box1.text_frame
    tf1.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf1.word_wrap = True
    pp1 = tf1.paragraphs[0]; pp1.alignment = PP_ALIGN.CENTER
    rr1 = pp1.add_run(); rr1.text = p["structure"]
    rr1.font.size = Pt(26); rr1.font.bold = True; rr1.font.color.rgb = WHITE
    # 中文
    textbox(s, Inches(0.7), Inches(2.65), Inches(11.93), Inches(0.45),
            [(p["zh"], 20, True, DARK, PP_ALIGN.CENTER)])
    # 例句標題
    textbox(s, Inches(0.7), Inches(3.20), Inches(11.5), Inches(0.4),
            [("例句 Examples", 18, True, THEME_D, PP_ALIGN.LEFT)])
    # 例句卡片
    ey = 3.62
    for en, zh in p["examples"]:
        eb = rrect(s, Inches(0.7), Inches(ey), Inches(12.0), Inches(0.48),
                   WHITE, line=CARD_BD, line_w=1.0)
        tf = eb.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        pp = tf.paragraphs[0]; pp.alignment = PP_ALIGN.LEFT
        r1 = pp.add_run(); r1.text = "  • " + en
        r1.font.size = Pt(17); r1.font.bold = True; r1.font.color.rgb = THEME_D
        r2 = pp.add_run(); r2.text = "       " + zh
        r2.font.size = Pt(15); r2.font.color.rgb = GRAY
        ey += 0.54
    # 填空標題
    textbox(s, Inches(0.7), Inches(4.78), Inches(11.5), Inches(0.4),
            [("填空練習 Fill in the blank", 18, True, THEME_D, PP_ALIGN.LEFT)])
    # 填空題
    fy = 5.2
    letters = ["A", "B", "C"]
    for j, item in enumerate(p["fill"]):
        qb = rrect(s, Inches(0.7), Inches(fy), Inches(12.0), Inches(0.5),
                   THEME, line=None)
        tfq = qb.text_frame
        tfq.vertical_anchor = MSO_ANCHOR.MIDDLE
        tfq.word_wrap = True
        pq = tfq.paragraphs[0]; pq.alignment = PP_ALIGN.LEFT
        rq = pq.add_run(); rq.text = "  %d.  %s" % (j + 1, item["q"])
        rq.font.size = Pt(17); rq.font.bold = True; rq.font.color.rgb = WHITE
        # 三選項橫排
        oy = fy + 0.55
        opt_gap = 0.1
        ow_in = (12.0 - 2 * opt_gap) / 3
        for k, opt in enumerate(item["opts"]):
            ox = Inches(0.7 + k * (ow_in + opt_gap))
            ob = rrect(s, ox, Inches(oy), Inches(ow_in), Inches(0.45),
                       WHITE, line=CARD_BD, line_w=1.0)
            tfo = ob.text_frame
            tfo.vertical_anchor = MSO_ANCHOR.MIDDLE
            tfo.word_wrap = True
            po = tfo.paragraphs[0]; po.alignment = PP_ALIGN.CENTER
            ro = po.add_run(); ro.text = "(%s) %s" % (letters[k], opt)
            ro.font.size = Pt(17); ro.font.color.rgb = DARK
        fy += 1.1


def s_dialogue(prs):
    s = slide_blank(prs)
    set_bg(s, LIGHT)
    topbar(s, "Dialogue 會話")
    heading(s, "Ken 和 Sue 的對話")
    y = 1.7
    bw = Inches(8.3)
    for role, name, en, zh in DIALOGUE:
        if role == "A":
            x = Inches(0.5)
            col = A_BLUE
        else:
            x = Inches(13.333 - 0.5) - bw
            col = B_ORNG
        bub = rrect(s, x, Inches(y), bw, Inches(0.86), col, line=None)
        tf = bub.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        p1 = tf.paragraphs[0]; p1.alignment = PP_ALIGN.LEFT
        r1 = p1.add_run(); r1.text = "  %s: %s" % (name, en)
        r1.font.size = Pt(20); r1.font.bold = True; r1.font.color.rgb = WHITE
        r2 = p1.add_run(); r2.text = "   " + zh
        r2.font.size = Pt(16); r2.font.color.rgb = RGBColor(0xF2, 0xF2, 0xF2)
        y += 0.95


def s_answers(prs):
    s = slide_blank(prs)
    set_bg(s, LIGHT)
    topbar(s, "Answers 解答")
    heading(s, "練習解答")
    letters = ["A", "B", "C"]
    rows = [
        ("連連看一",         MATCH1_ANS),
        ("連連看二",         MATCH2_ANS),
        ("句型 1 · 第 1 題", "✓ (%s) %s" % (
            letters[PATTERN1["fill"][0]["ans"]],
            PATTERN1["fill"][0]["opts"][PATTERN1["fill"][0]["ans"]])),
        ("句型 1 · 第 2 題", "✓ (%s) %s" % (
            letters[PATTERN1["fill"][1]["ans"]],
            PATTERN1["fill"][1]["opts"][PATTERN1["fill"][1]["ans"]])),
        ("句型 2 · 第 1 題", "✓ (%s) %s" % (
            letters[PATTERN2["fill"][0]["ans"]],
            PATTERN2["fill"][0]["opts"][PATTERN2["fill"][0]["ans"]])),
        ("句型 2 · 第 2 題", "✓ (%s) %s" % (
            letters[PATTERN2["fill"][1]["ans"]],
            PATTERN2["fill"][1]["opts"][PATTERN2["fill"][1]["ans"]])),
    ]
    y = 1.7
    row_h = 0.78
    gap = 0.12
    for label, ans in rows:
        ab = rrect(s, Inches(0.7), Inches(y), Inches(11.93), Inches(row_h),
                   WHITE, line=GREEN, line_w=1.5)
        tf = ab.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        pp = tf.paragraphs[0]; pp.alignment = PP_ALIGN.LEFT
        r1 = pp.add_run(); r1.text = "  " + label + "   "
        r1.font.size = Pt(20); r1.font.bold = True; r1.font.color.rgb = DARK
        r2 = pp.add_run(); r2.text = ans
        r2.font.size = Pt(20); r2.font.bold = True; r2.font.color.rgb = GREEN
        y += row_h + gap


def s_closing(prs):
    s = slide_blank(prs)
    set_bg(s, THEME)
    textbox(s, Inches(0.8), Inches(2.2), Inches(11.7), Inches(1.6),
            [("🎉", 90, False, WHITE, PP_ALIGN.CENTER)])
    textbox(s, Inches(0.8), Inches(3.8), Inches(11.7), Inches(1.1),
            [("Great Job!  你做得很棒！", 46, True, WHITE, PP_ALIGN.CENTER)])
    textbox(s, Inches(0.8), Inches(4.95), Inches(11.7), Inches(0.8),
            [("Unit %d · %s  %s" % (UNIT_NO, UNIT_EN, UNIT_ZH),
              22, False, WHITE, PP_ALIGN.CENTER)])
    textbox(s, Inches(0.8), Inches(5.85), Inches(11.7), Inches(0.7),
            [("🌐 線上互動版　https://Yuchiaoniu.github.io/english/",
              20, False, WHITE, PP_ALIGN.CENTER)])


def main():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    s_cover(prs)                                                                # 1
    s_vocab_overview(prs)                                                       # 2
    s_vocab_cards(prs, VOCAB[0:3], "Vocabulary 1/3", "情緒詞（一）")           # 3
    s_vocab_cards(prs, VOCAB[3:6], "Vocabulary 2/3", "情緒詞（二）")           # 4
    s_vocab_cards(prs, VOCAB[6:9], "Vocabulary 3/3", "情緒詞（三）")           # 5
    s_section(prs, "🎯", "Matching Game", "連連看")                            # 6
    s_match(prs, "Matching 1 連連看（一）",
            "連連看（一）：單字 ↔ 中文",
            "請將左邊英文單字與右邊中文意義配對。",
            MATCH1_LEFT, MATCH1_RIGHT)                                          # 7
    s_match(prs, "Matching 2 連連看（二）",
            "連連看（二）：情境 ↔ 情緒",
            "請看左邊情境句，從右邊找出對應的情緒詞。",
            MATCH2_LEFT, MATCH2_RIGHT)                                          # 8
    s_section(prs, "📝", "Pattern & Fill-in", "句型填空")                      # 9
    s_pattern_fill(prs, 1, PATTERN1)                                            # 10
    s_pattern_fill(prs, 2, PATTERN2)                                            # 11
    s_section(prs, "💬", "Dialogue", "會話")                                   # 12
    s_dialogue(prs)                                                             # 13
    s_answers(prs)                                                              # 14
    s_closing(prs)                                                              # 15

    out = r"C:\Users\yuchi\grade4-english-tutor\Unit9_How_Do_You_Feel.pptx"
    prs.save(out)
    print("Saved: %s  (%d slides)" % (out, len(prs.slides._sldIdLst)))


if __name__ == "__main__":
    main()
