# -*- coding: utf-8 -*-
"""Unit 7 — What Time Is It? 單課完整簡報產生器。
涵蓋：封面、單字(16)、句型(2)、對話(6)、測驗(5)+解答、大學伴提示、結尾。
設計目標：15~20 頁、配色精美、字體大。
輸出：Unit7_What_Time_Is_It.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------- 色彩 ----------
THEME   = RGBColor(0x16, 0xA0, 0x85)   # 第七課主色（teal）
THEME_D = RGBColor(0x0E, 0x6B, 0x59)   # 深一階
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
DARK    = RGBColor(0x2C, 0x3E, 0x50)
GRAY    = RGBColor(0x7F, 0x8C, 0x8D)
LIGHT   = RGBColor(0xF4, 0xF8, 0xF7)
A_BLUE  = RGBColor(0x4A, 0x90, 0xD9)   # 對話 A
B_ORNG  = RGBColor(0xF0, 0x8A, 0x24)   # 對話 B
GREEN   = RGBColor(0x27, 0xAE, 0x60)   # 正解
CARD_BD = RGBColor(0xD8, 0xE6, 0xE2)

W = Inches(13.333)
H = Inches(7.5)

# ---------- 第七課資料 ----------
# 數字採用「鐘面」emoji（🕐..🕙），在 PowerPoint 可彩色渲染且呼應報時主題；
# keycap emoji（1️⃣..🔟）在 Segoe UI Emoji 放大後會變空白，故不用於簡報。
NUMBERS = [
    {"word": "one",   "zh": "一", "zhuyin": "ㄧ",      "emoji": "🕐"},
    {"word": "two",   "zh": "二", "zhuyin": "ㄦˋ",     "emoji": "🕑"},
    {"word": "three", "zh": "三", "zhuyin": "ㄙㄢ",    "emoji": "🕒"},
    {"word": "four",  "zh": "四", "zhuyin": "ㄙˋ",     "emoji": "🕓"},
    {"word": "five",  "zh": "五", "zhuyin": "ㄨˇ",     "emoji": "🕔"},
    {"word": "six",   "zh": "六", "zhuyin": "ㄌㄧㄡˋ", "emoji": "🕕"},
    {"word": "seven", "zh": "七", "zhuyin": "ㄑㄧ",    "emoji": "🕖"},
    {"word": "eight", "zh": "八", "zhuyin": "ㄅㄚ",    "emoji": "🕗"},
    {"word": "nine",  "zh": "九", "zhuyin": "ㄐㄧㄡˇ", "emoji": "🕘"},
    {"word": "ten",   "zh": "十", "zhuyin": "ㄕˊ",     "emoji": "🕙"},
]
TIMEWORDS = [
    {"word": "clock",     "zh": "時鐘",    "zhuyin": "ㄕˊ ㄓㄨㄥ",   "emoji": "🕐"},
    {"word": "time",      "zh": "時間",    "zhuyin": "ㄕˊ ㄐㄧㄢ",   "emoji": "⌚"},
    {"word": "morning",   "zh": "早上",    "zhuyin": "ㄗㄠˇ ˙ㄕㄤ",  "emoji": "🌅"},
    {"word": "afternoon", "zh": "下午",    "zhuyin": "ㄒㄧㄚˋ ˙ㄨ",  "emoji": "☀️"},
    {"word": "night",     "zh": "晚上",    "zhuyin": "ㄨㄢˇ ˙ㄕㄤ",  "emoji": "🌙"},
    {"word": "o'clock",   "zh": "...點鐘", "zhuyin": "ㄉㄧㄢˇ ㄓㄨㄥ", "emoji": "⏰"},
]
ALL_VOCAB = NUMBERS + TIMEWORDS

PATTERNS = [
    {
        "structure": "What time is it?  It's ___ o'clock.",
        "zh": "現在幾點？　現在 ___ 點。",
        "examples": [
            ("It's three o'clock.", "現在三點。"),
            ("It's seven o'clock.", "現在七點。"),
        ],
    },
    {
        "structure": "I ___ at ___ o'clock.",
        "zh": "我 ___ 點 ___。（描述每天作息）",
        "examples": [
            ("I eat breakfast at seven o'clock.", "我七點吃早餐。"),
            ("I go to school at eight o'clock.", "我八點上學。"),
        ],
    },
]

DIALOGUE = [
    ("A", "Ken", "Excuse me. What time is it?", "打擾一下。現在幾點？"),
    ("B", "Sue", "It's three o'clock in the afternoon.", "現在下午三點。"),
    ("A", "Ken", "Oh no! I'm late for class!", "糟糕！我上課要遲到了！"),
    ("B", "Sue", "Class starts at three, right? Run!", "三點上課，對嗎？快跑！"),
    ("A", "Ken", "When does class end?", "課什麼時候結束？"),
    ("B", "Sue", "At five o'clock. Don't worry!", "五點。別擔心！"),
]

QUIZ = [
    {"q": "「It's seven o'clock.」是什麼意思？",
     "opts": ["現在七點了", "我七歲了", "有七個人", "第七天"], "ans": 0},
    {"q": "下午的英文是？",
     "opts": ["morning", "night", "afternoon", "evening"], "ans": 2},
    {"q": "怎麼問現在幾點？",
     "opts": ["Where is the clock?", "What time is it?", "How many clocks?", "Is it time?"], "ans": 1},
    {"q": "對話中課程幾點結束？",
     "opts": ["三點", "四點", "五點", "六點"], "ans": 2},
    {"q": "Which word means 早上?",
     "opts": ["night", "afternoon", "o'clock", "morning"], "ans": 3},
]

TUTOR = {
    "key": [
        "用時鐘玩具或畫時鐘練習整點報時。",
        "morning / afternoon / night 對應一天的三個時段。",
        "o'clock 拼字有撇號 ( ' )，要帶學生練習寫。",
    ],
    "errors": [
        "three 的 th 發 /θ/，不是 /s/。",
        "It's three. 容易漏掉 o'clock，應說 It's three o'clock.",
        "What time is it? 常被簡化成 What time?",
    ],
    "prompts": [
        "用 It's ___ o'clock 報出現在時間。",
        "用 I ___ at ___ o'clock 介紹你的一天作息。",
        "你幾點起床？幾點睡覺？最喜歡哪個時段？",
    ],
}

# ---------- 共用工具 ----------
def slide_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def textbox(slide, x, y, w, h, lines, anchor=MSO_ANCHOR.MIDDLE, wrap=True):
    """lines: list of (text, size, bold, color, align)."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for i, (text, size, bold, color, align) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return box


def rrect(slide, x, y, w, h, fill, line=None, line_w=1.0):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.shadow.inherit = False
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
    return sh


def topbar(slide, label):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.62))
    bar.shadow.inherit = False
    bar.fill.solid()
    bar.fill.fore_color.rgb = THEME
    bar.line.fill.background()
    tf = bar.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "   Unit 7 · What Time Is It?   ⏰   " + label
    r.font.size = Pt(17)
    r.font.bold = True
    r.font.color.rgb = WHITE


def heading(slide, text):
    textbox(slide, Inches(0.5), Inches(0.72), Inches(12.3), Inches(0.85),
            [(text, 32, True, THEME_D, PP_ALIGN.LEFT)], anchor=MSO_ANCHOR.MIDDLE)


# ---------- 投影片 ----------
def s_cover(prs):
    s = slide_blank(prs)
    set_bg(s, THEME)
    # 裝飾圓
    deco = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.3), Inches(-1.4),
                              Inches(4.6), Inches(4.6))
    deco.shadow.inherit = False
    deco.fill.solid(); deco.fill.fore_color.rgb = THEME_D
    deco.line.fill.background()
    deco2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.6), Inches(4.6),
                               Inches(4.2), Inches(4.2))
    deco2.shadow.inherit = False
    deco2.fill.solid(); deco2.fill.fore_color.rgb = THEME_D
    deco2.line.fill.background()

    textbox(s, Inches(0.8), Inches(1.4), Inches(11.7), Inches(1.7),
            [("⏰", 96, False, WHITE, PP_ALIGN.CENTER)])
    textbox(s, Inches(0.8), Inches(3.05), Inches(11.7), Inches(1.2),
            [("Unit 7 · What Time Is It?", 50, True, WHITE, PP_ALIGN.CENTER)])
    textbox(s, Inches(0.8), Inches(4.25), Inches(11.7), Inches(0.9),
            [("現在幾點？", 40, True, WHITE, PP_ALIGN.CENTER)])
    textbox(s, Inches(0.8), Inches(5.45), Inches(11.7), Inches(0.7),
            [("四年級英語 · 單字 16 · 句型 2 · 對話 · 測驗", 22, False, WHITE, PP_ALIGN.CENTER)])
    textbox(s, Inches(0.8), Inches(6.35), Inches(11.7), Inches(0.55),
            [("🌐  https://Yuchiaoniu.github.io/english/", 18, False, WHITE, PP_ALIGN.CENTER)])


def s_vocab_overview(prs):
    s = slide_blank(prs)
    set_bg(s, LIGHT)
    topbar(s, "Vocabulary 單字總覽")
    heading(s, "本課單字 16 個")
    cols, rows = 4, 4
    cw, ch = Inches(3.05), Inches(1.18)
    gx, gy = Inches(0.18), Inches(0.12)
    sx, sy = Inches(0.36), Inches(1.65)
    for i, v in enumerate(ALL_VOCAB):
        c, r = i % cols, i // cols
        x = sx + c * (cw + gx)
        y = sy + r * (ch + gy)
        card = rrect(s, x, y, cw, ch, WHITE, line=CARD_BD, line_w=1.25)
        tf = card.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        p1 = tf.paragraphs[0]; p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run(); r1.text = "%s  %s" % (v["emoji"], v["word"])
        r1.font.size = Pt(20); r1.font.bold = True; r1.font.color.rgb = DARK
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = v["zh"]
        r2.font.size = Pt(16); r2.font.color.rgb = THEME


def s_vocab_cards(prs, items, label, title):
    s = slide_blank(prs)
    set_bg(s, LIGHT)
    topbar(s, label)
    heading(s, title)
    n = len(items)
    usable = 12.5
    cw = (usable - 0.3 * (n - 1)) / n
    cwI = Inches(cw)
    chI = Inches(4.3)
    sx = Inches((13.333 - (cw * n + 0.3 * (n - 1))) / 2)
    sy = Inches(1.95)
    for i, v in enumerate(items):
        x = Emu(int(sx) + i * (int(cwI) + int(Inches(0.3))))
        card = rrect(s, x, sy, cwI, chI, WHITE, line=THEME, line_w=1.75)
        tf = card.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        p0 = tf.paragraphs[0]; p0.alignment = PP_ALIGN.CENTER
        r0 = p0.add_run(); r0.text = v["emoji"]; r0.font.size = Pt(60)
        p1 = tf.add_paragraph(); p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run(); r1.text = v["word"]
        r1.font.size = Pt(30); r1.font.bold = True; r1.font.color.rgb = THEME_D
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = v["zh"]
        r2.font.size = Pt(24); r2.font.bold = True; r2.font.color.rgb = DARK
        p3 = tf.add_paragraph(); p3.alignment = PP_ALIGN.CENTER
        r3 = p3.add_run(); r3.text = v["zhuyin"]
        r3.font.size = Pt(18); r3.font.color.rgb = GRAY


def s_section(prs, emoji, en, zh):
    s = slide_blank(prs)
    set_bg(s, THEME)
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.55), W, Inches(2.4))
    band.shadow.inherit = False
    band.fill.solid(); band.fill.fore_color.rgb = THEME_D
    band.line.fill.background()
    textbox(s, Inches(0.8), Inches(1.15), Inches(11.7), Inches(1.4),
            [(emoji, 80, False, WHITE, PP_ALIGN.CENTER)])
    textbox(s, Inches(0.8), Inches(2.75), Inches(11.7), Inches(1.1),
            [(en, 48, True, WHITE, PP_ALIGN.CENTER)])
    textbox(s, Inches(0.8), Inches(3.9), Inches(11.7), Inches(0.9),
            [(zh, 32, False, WHITE, PP_ALIGN.CENTER)])


def s_pattern(prs, idx, p):
    s = slide_blank(prs)
    set_bg(s, LIGHT)
    topbar(s, "Sentence Pattern %d/2 句型" % idx)
    heading(s, "句型 %d" % idx)
    # 句型結構框
    box = rrect(s, Inches(0.7), Inches(1.75), Inches(11.93), Inches(1.5),
                THEME, line=None)
    tf = box.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.word_wrap = True
    pp = tf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
    rr = pp.add_run(); rr.text = p["structure"]
    rr.font.size = Pt(34); rr.font.bold = True; rr.font.color.rgb = WHITE
    # 中文說明
    textbox(s, Inches(0.7), Inches(3.45), Inches(11.93), Inches(0.8),
            [(p["zh"], 26, True, DARK, PP_ALIGN.CENTER)])
    # 例句
    textbox(s, Inches(0.9), Inches(4.4), Inches(11.5), Inches(0.55),
            [("例句 Examples", 22, True, THEME_D, PP_ALIGN.LEFT)])
    ey = 5.0
    for en, zh in p["examples"]:
        eb = rrect(s, Inches(0.9), Inches(ey), Inches(11.5), Inches(0.95),
                   WHITE, line=CARD_BD, line_w=1.25)
        tf2 = eb.text_frame; tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf2.word_wrap = True
        p1 = tf2.paragraphs[0]; p1.alignment = PP_ALIGN.LEFT
        r1 = p1.add_run(); r1.text = "  " + en
        r1.font.size = Pt(24); r1.font.bold = True; r1.font.color.rgb = THEME_D
        r2 = p1.add_run(); r2.text = "    " + zh
        r2.font.size = Pt(20); r2.font.color.rgb = GRAY
        ey += 1.05


def s_dialogue(prs):
    s = slide_blank(prs)
    set_bg(s, LIGHT)
    topbar(s, "Dialogue 對話")
    heading(s, "對話：Ken 和 Sue")
    y = 1.7
    bw = Inches(8.3)
    for role, name, en, zh in DIALOGUE:
        if role == "A":
            x = Inches(0.5); col = A_BLUE
        else:
            x = Inches(13.333 - 0.5) - bw; col = B_ORNG
        bub = rrect(s, x, Inches(y), bw, Inches(0.86), col, line=None)
        tf = bub.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        p1 = tf.paragraphs[0]; p1.alignment = PP_ALIGN.LEFT
        r1 = p1.add_run(); r1.text = "  %s: %s" % (name, en)
        r1.font.size = Pt(21); r1.font.bold = True; r1.font.color.rgb = WHITE
        r2 = p1.add_run(); r2.text = "   %s" % zh
        r2.font.size = Pt(17); r2.font.color.rgb = RGBColor(0xF2, 0xF2, 0xF2)
        y += 0.95


def s_quiz(prs, items, start, label):
    s = slide_blank(prs)
    set_bg(s, LIGHT)
    topbar(s, label)
    heading(s, "小測驗 Quiz")
    letters = ["A", "B", "C", "D"]
    bar_h, opt_h, row_gap, blk_gap = 0.6, 0.48, 0.05, 0.18
    y = 1.7
    for j, item in enumerate(items):
        qno = start + j
        qb = rrect(s, Inches(0.6), Inches(y), Inches(12.13), Inches(bar_h),
                   THEME, line=None)
        tf = qb.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        pp = tf.paragraphs[0]; pp.alignment = PP_ALIGN.LEFT
        rr = pp.add_run(); rr.text = "  Q%d.  %s" % (qno, item["q"])
        rr.font.size = Pt(21); rr.font.bold = True; rr.font.color.rgb = WHITE
        # 選項 2x2
        oy = y + bar_h + row_gap
        ow = Inches(5.95)
        for k, opt in enumerate(item["opts"]):
            ox = Inches(0.6) if k % 2 == 0 else Inches(6.78)
            row = k // 2
            oyy = oy + row * (opt_h + row_gap)
            ob = rrect(s, ox, Inches(oyy), ow, Inches(opt_h),
                       WHITE, line=CARD_BD, line_w=1.0)
            tf2 = ob.text_frame; tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.LEFT
            r2 = p2.add_run(); r2.text = "  (%s) %s" % (letters[k], opt)
            r2.font.size = Pt(18); r2.font.color.rgb = DARK
        y = oy + 2 * opt_h + row_gap + blk_gap


def s_answers(prs):
    s = slide_blank(prs)
    set_bg(s, LIGHT)
    topbar(s, "Answers 解答")
    heading(s, "測驗解答")
    letters = ["A", "B", "C", "D"]
    y = 1.8
    for i, item in enumerate(QUIZ):
        ab = rrect(s, Inches(1.1), Inches(y), Inches(11.1), Inches(0.86),
                   WHITE, line=GREEN, line_w=1.75)
        tf = ab.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        pp = tf.paragraphs[0]; pp.alignment = PP_ALIGN.LEFT
        r0 = pp.add_run(); r0.text = "  Q%d   " % (i + 1)
        r0.font.size = Pt(22); r0.font.bold = True; r0.font.color.rgb = DARK
        r1 = pp.add_run(); r1.text = "✓ (%s) %s" % (letters[item["ans"]], item["opts"][item["ans"]])
        r1.font.size = Pt(22); r1.font.bold = True; r1.font.color.rgb = GREEN
        y += 1.0


def s_tutor(prs):
    s = slide_blank(prs)
    set_bg(s, LIGHT)
    topbar(s, "Tutor Guide 大學伴提示")
    heading(s, "給大學伴的教學提示")
    blocks = [
        ("🎯 教學重點", TUTOR["key"], THEME),
        ("⚠️ 常見錯誤", TUTOR["errors"], B_ORNG),
        ("💬 互動引導語", TUTOR["prompts"], A_BLUE),
    ]
    y = 1.65
    for title, items, col in blocks:
        hb = rrect(s, Inches(0.6), Inches(y), Inches(12.13), Inches(0.55),
                   col, line=None)
        tf = hb.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        pp = tf.paragraphs[0]; pp.alignment = PP_ALIGN.LEFT
        rr = pp.add_run(); rr.text = "  " + title
        rr.font.size = Pt(22); rr.font.bold = True; rr.font.color.rgb = WHITE
        bb = rrect(s, Inches(0.6), Inches(y + 0.6), Inches(12.13), Inches(1.18),
                   WHITE, line=CARD_BD, line_w=1.0)
        tf2 = bb.text_frame; tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf2.word_wrap = True
        for k, it in enumerate(items):
            p = tf2.paragraphs[0] if k == 0 else tf2.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run(); r.text = "•  " + it
            r.font.size = Pt(18); r.font.color.rgb = DARK
        y += 1.92


def s_closing(prs):
    s = slide_blank(prs)
    set_bg(s, THEME)
    textbox(s, Inches(0.8), Inches(2.2), Inches(11.7), Inches(1.6),
            [("🎉", 90, False, WHITE, PP_ALIGN.CENTER)])
    textbox(s, Inches(0.8), Inches(3.8), Inches(11.7), Inches(1.1),
            [("Good Job!  你做得很棒！", 46, True, WHITE, PP_ALIGN.CENTER)])
    textbox(s, Inches(0.8), Inches(4.95), Inches(11.7), Inches(0.8),
            [("Unit 7 · What Time Is It?  現在幾點？", 24, False, WHITE, PP_ALIGN.CENTER)])
    textbox(s, Inches(0.8), Inches(5.85), Inches(11.7), Inches(0.7),
            [("🌐 線上互動版　https://Yuchiaoniu.github.io/english/",
              20, False, WHITE, PP_ALIGN.CENTER)])


def main():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    s_cover(prs)                                            # 1
    s_vocab_overview(prs)                                   # 2
    s_vocab_cards(prs, NUMBERS[0:5], "Numbers 1-5 數字", "數字 1–5")    # 3
    s_vocab_cards(prs, NUMBERS[5:10], "Numbers 6-10 數字", "數字 6–10")  # 4
    s_vocab_cards(prs, TIMEWORDS[0:3], "Time Words 時間相關字", "時間相關字（一）")  # 5
    s_vocab_cards(prs, TIMEWORDS[3:6], "Time Words 時間相關字", "時間相關字（二）")  # 6
    s_section(prs, "📝", "Sentence Patterns", "句型練習")    # 7
    s_pattern(prs, 1, PATTERNS[0])                          # 8
    s_pattern(prs, 2, PATTERNS[1])                          # 9
    s_section(prs, "💬", "Dialogue", "對話閱讀")            # 10
    s_dialogue(prs)                                         # 11
    s_section(prs, "✅", "Quiz", "單元測驗")                # 12
    s_quiz(prs, QUIZ[0:3], 1, "Quiz 1-3 測驗")              # 13
    s_quiz(prs, QUIZ[3:5], 4, "Quiz 4-5 測驗")              # 14
    s_answers(prs)                                          # 15
    s_tutor(prs)                                            # 16
    s_closing(prs)                                          # 17

    out = r"C:\Users\yuchi\grade4-english-tutor\Unit7_What_Time_Is_It.pptx"
    prs.save(out)
    print("Saved: %s  (%d slides)" % (out, len(prs.slides._sldIdLst)))


if __name__ == "__main__":
    main()
