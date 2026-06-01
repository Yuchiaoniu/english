from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

UNITS = [
    {
        "id": 1, "title": "Hello, Friends!", "emoji": "👋",
        "vocab": [
            {"word": "hello",    "zh": "你好",   "zhuyin": "ㄋㄧˇ ㄏㄠˇ",       "emoji": "👋"},
            {"word": "goodbye",  "zh": "再見",   "zhuyin": "ㄗㄞˋ ㄐㄧㄢˋ",     "emoji": "🙋"},
            {"word": "friend",   "zh": "朋友",   "zhuyin": "ㄆㄥˊ ㄧㄡˇ",       "emoji": "👫"},
            {"word": "name",     "zh": "名字",   "zhuyin": "ㄇㄧㄥˊ ˙ㄗ",       "emoji": "🏷️"},
            {"word": "class",    "zh": "班級",   "zhuyin": "ㄅㄢ ㄐㄧˊ",        "emoji": "🎒"},
            {"word": "teacher",  "zh": "老師",   "zhuyin": "ㄌㄠˇ ㄕ",          "emoji": "👩‍🏫"},
            {"word": "student",  "zh": "學生",   "zhuyin": "ㄒㄩㄝˊ ㄕㄥ",       "emoji": "🧒"},
            {"word": "boy",      "zh": "男生",   "zhuyin": "ㄋㄢˊ ㄕㄥ",        "emoji": "👦"},
            {"word": "girl",     "zh": "女生",   "zhuyin": "ㄋㄩˇ ㄕㄥ",        "emoji": "👧"},
            {"word": "nice",     "zh": "好的/棒的","zhuyin": "ㄏㄠˇ",           "emoji": "😊"},
            {"word": "meet",     "zh": "認識",   "zhuyin": "ㄖㄣˋ ˙ㄕ",        "emoji": "🤝"},
            {"word": "new",      "zh": "新的",   "zhuyin": "ㄒㄧㄣ",            "emoji": "✨"},
            {"word": "please",   "zh": "請",     "zhuyin": "ㄑㄧㄥˇ",           "emoji": "🙏"},
            {"word": "yes",      "zh": "是/對",  "zhuyin": "ㄉㄨㄟˋ",           "emoji": "✅"},
            {"word": "no",       "zh": "不",     "zhuyin": "ㄅㄨˋ",             "emoji": "❌"},
            {"word": "sorry",    "zh": "對不起", "zhuyin": "ㄉㄨㄟˋ ˙ㄅㄨ ㄑㄧˇ", "emoji": "😔"},
        ],
        "color": RGBColor(0x4A, 0x90, 0xD9),
    },
    {
        "id": 2, "title": "My Family", "emoji": "👨‍👩‍👧‍👦",
        "vocab": [
            {"word": "family",   "zh": "家人",     "zhuyin": "ㄐㄧㄚ ㄖㄣˊ",       "emoji": "👨‍👩‍👧‍👦"},
            {"word": "father",   "zh": "爸爸",     "zhuyin": "ㄅㄚˋ ˙ㄅㄚ",        "emoji": "👨"},
            {"word": "mother",   "zh": "媽媽",     "zhuyin": "ㄇㄚ ˙ㄇㄚ",         "emoji": "👩"},
            {"word": "brother",  "zh": "哥哥/弟弟", "zhuyin": "ㄍㄜ ˙ㄍㄜ",         "emoji": "👦"},
            {"word": "sister",   "zh": "姐姐/妹妹", "zhuyin": "ㄐㄧㄝˇ ˙ㄐㄧㄝ",    "emoji": "👧"},
            {"word": "grandpa",  "zh": "爺爺",     "zhuyin": "ㄧㄝˊ ˙ㄧㄝ",        "emoji": "👴"},
            {"word": "grandma",  "zh": "奶奶",     "zhuyin": "ㄋㄞˇ ˙ㄋㄞ",        "emoji": "👵"},
            {"word": "baby",     "zh": "嬰兒",     "zhuyin": "ㄧㄥ ㄦˊ",           "emoji": "👶"},
            {"word": "uncle",    "zh": "叔叔/伯伯", "zhuyin": "ㄕㄨˊ ˙ㄕㄨ",       "emoji": "🧔"},
            {"word": "aunt",     "zh": "阿姨",     "zhuyin": "ㄚ ㄧˊ",             "emoji": "👩"},
            {"word": "cousin",   "zh": "堂/表兄弟姐妹","zhuyin": "ㄅㄧㄠˇ ㄒㄩㄥ",  "emoji": "👫"},
            {"word": "son",      "zh": "兒子",     "zhuyin": "ㄦˊ ˙ㄗ",            "emoji": "👦"},
            {"word": "daughter", "zh": "女兒",     "zhuyin": "ㄋㄩˇ ㄦˊ",          "emoji": "👧"},
            {"word": "home",     "zh": "家",       "zhuyin": "ㄐㄧㄚ",             "emoji": "🏠"},
            {"word": "love",     "zh": "愛",       "zhuyin": "ㄞˋ",               "emoji": "❤️"},
            {"word": "happy",    "zh": "快樂的",   "zhuyin": "ㄎㄨㄞˋ ㄌㄜˋ",      "emoji": "😄"},
        ],
        "color": RGBColor(0xE8, 0x74, 0x6A),
    },
    {
        "id": 3, "title": "At School", "emoji": "🏫",
        "vocab": [
            {"word": "school",    "zh": "學校",   "zhuyin": "ㄒㄩㄝˊ ㄒㄧㄠˋ",  "emoji": "🏫"},
            {"word": "classroom", "zh": "教室",   "zhuyin": "ㄐㄧㄠˋ ㄕˋ",      "emoji": "🚪"},
            {"word": "desk",      "zh": "桌子",   "zhuyin": "ㄓㄨㄛ ˙ㄗ",       "emoji": "🪑"},
            {"word": "chair",     "zh": "椅子",   "zhuyin": "ㄧˇ ˙ㄗ",          "emoji": "💺"},
            {"word": "book",      "zh": "書",     "zhuyin": "ㄕㄨ",             "emoji": "📚"},
            {"word": "pencil",    "zh": "鉛筆",   "zhuyin": "ㄑㄧㄢ ㄅㄧˇ",     "emoji": "✏️"},
            {"word": "ruler",     "zh": "尺",     "zhuyin": "ㄔˇ",              "emoji": "📏"},
            {"word": "bag",       "zh": "書包",   "zhuyin": "ㄕㄨ ㄅㄠ",        "emoji": "🎒"},
            {"word": "eraser",    "zh": "橡皮擦", "zhuyin": "ㄒㄧㄤˋ ㄆㄧˊ ㄘㄚ","emoji": "🧹"},
            {"word": "crayon",    "zh": "蠟筆",   "zhuyin": "ㄌㄚˋ ㄅㄧˇ",      "emoji": "🖍️"},
            {"word": "notebook",  "zh": "筆記本", "zhuyin": "ㄅㄧˇ ㄐㄧˋ ˙ㄅㄣ",  "emoji": "📓"},
            {"word": "pen",       "zh": "原子筆", "zhuyin": "ㄩㄢˊ ˙ㄗ ㄅㄧˇ",   "emoji": "🖊️"},
            {"word": "scissors",  "zh": "剪刀",   "zhuyin": "ㄐㄧㄢˇ ㄉㄠ",      "emoji": "✂️"},
            {"word": "glue",      "zh": "膠水",   "zhuyin": "ㄐㄧㄠ ㄕㄨㄟˇ",   "emoji": "🧴"},
            {"word": "board",     "zh": "黑板",   "zhuyin": "ㄏㄟ ㄅㄢˇ",       "emoji": "🖼️"},
            {"word": "door",      "zh": "門",     "zhuyin": "ㄇㄣˊ",            "emoji": "🚪"},
        ],
        "color": RGBColor(0x5B, 0xB5, 0x7E),
    },
    {
        "id": 4, "title": "I'm Hungry!", "emoji": "🍽️",
        "vocab": [
            {"word": "apple",    "zh": "蘋果",   "zhuyin": "ㄆㄧㄥˊ ㄍㄨㄛˇ",  "emoji": "🍎"},
            {"word": "banana",   "zh": "香蕉",   "zhuyin": "ㄒㄧㄤ ㄐㄧㄠ",    "emoji": "🍌"},
            {"word": "rice",     "zh": "飯",     "zhuyin": "ㄈㄢˋ",            "emoji": "🍚"},
            {"word": "milk",     "zh": "牛奶",   "zhuyin": "ㄋㄧㄡˊ ㄋㄞˇ",    "emoji": "🥛"},
            {"word": "juice",    "zh": "果汁",   "zhuyin": "ㄍㄨㄛˇ ㄓ",       "emoji": "🧃"},
            {"word": "bread",    "zh": "麵包",   "zhuyin": "ㄇㄧㄢˋ ㄅㄠ",     "emoji": "🍞"},
            {"word": "egg",      "zh": "雞蛋",   "zhuyin": "ㄐㄧ ㄉㄢˋ",       "emoji": "🥚"},
            {"word": "water",    "zh": "水",     "zhuyin": "ㄕㄨㄟˇ",          "emoji": "💧"},
            {"word": "cookie",   "zh": "餅乾",   "zhuyin": "ㄅㄧㄥˇ ㄍㄢ",     "emoji": "🍪"},
            {"word": "cake",     "zh": "蛋糕",   "zhuyin": "ㄉㄢˋ ㄍㄠ",       "emoji": "🎂"},
            {"word": "soup",     "zh": "湯",     "zhuyin": "ㄊㄤ",             "emoji": "🍲"},
            {"word": "noodles",  "zh": "麵條",   "zhuyin": "ㄇㄧㄢˋ ㄊㄧㄠˊ",  "emoji": "🍜"},
            {"word": "orange",   "zh": "柳橙",   "zhuyin": "ㄌㄧㄡˇ ㄔㄥˊ",    "emoji": "🍊"},
            {"word": "grape",    "zh": "葡萄",   "zhuyin": "ㄆㄨˊ ˙ㄊㄠ",      "emoji": "🍇"},
            {"word": "sandwich", "zh": "三明治",  "zhuyin": "ㄙㄢ ㄇㄧㄥˊ ㄓˋ", "emoji": "🥪"},
            {"word": "hungry",   "zh": "餓的",   "zhuyin": "ㄜˋ",              "emoji": "😋"},
        ],
        "color": RGBColor(0xF5, 0xA6, 0x23),
    },
    {
        "id": 5, "title": "Animals I Like", "emoji": "🐾",
        "vocab": [
            {"word": "cat",       "zh": "貓",    "zhuyin": "ㄇㄠ",            "emoji": "🐱"},
            {"word": "dog",       "zh": "狗",    "zhuyin": "ㄍㄡˇ",           "emoji": "🐶"},
            {"word": "rabbit",    "zh": "兔子",  "zhuyin": "ㄊㄨˋ ˙ㄗ",      "emoji": "🐰"},
            {"word": "bird",      "zh": "鳥",    "zhuyin": "ㄋㄧㄠˇ",         "emoji": "🐦"},
            {"word": "fish",      "zh": "魚",    "zhuyin": "ㄩˊ",             "emoji": "🐟"},
            {"word": "bear",      "zh": "熊",    "zhuyin": "ㄒㄩㄥˊ",         "emoji": "🐻"},
            {"word": "tiger",     "zh": "老虎",  "zhuyin": "ㄌㄠˇ ㄏㄨˇ",    "emoji": "🐯"},
            {"word": "elephant",  "zh": "大象",  "zhuyin": "ㄉㄚˋ ㄒㄧㄤˋ",   "emoji": "🐘"},
            {"word": "monkey",    "zh": "猴子",  "zhuyin": "ㄏㄡˊ ˙ㄗ",      "emoji": "🐵"},
            {"word": "lion",      "zh": "獅子",  "zhuyin": "ㄕ ˙ㄗ",          "emoji": "🦁"},
            {"word": "panda",     "zh": "熊貓",  "zhuyin": "ㄒㄩㄥˊ ㄇㄠ",    "emoji": "🐼"},
            {"word": "horse",     "zh": "馬",    "zhuyin": "ㄇㄚˇ",           "emoji": "🐴"},
            {"word": "frog",      "zh": "青蛙",  "zhuyin": "ㄑㄧㄥ ㄨㄚ",     "emoji": "🐸"},
            {"word": "turtle",    "zh": "烏龜",  "zhuyin": "ㄨ ㄍㄨㄟ",       "emoji": "🐢"},
            {"word": "butterfly", "zh": "蝴蝶",  "zhuyin": "ㄏㄨˊ ㄉㄧㄝˊ",   "emoji": "🦋"},
            {"word": "duck",      "zh": "鴨子",  "zhuyin": "ㄧㄚ ˙ㄗ",        "emoji": "🦆"},
        ],
        "color": RGBColor(0x9B, 0x59, 0xB6),
    },
    {
        "id": 6, "title": "Dress Up!", "emoji": "👗",
        "vocab": [
            {"word": "red",      "zh": "紅色",  "zhuyin": "ㄏㄨㄥˊ ㄙㄜˋ",   "emoji": "🔴"},
            {"word": "blue",     "zh": "藍色",  "zhuyin": "ㄌㄢˊ ㄙㄜˋ",     "emoji": "🔵"},
            {"word": "yellow",   "zh": "黃色",  "zhuyin": "ㄏㄨㄤˊ ㄙㄜˋ",   "emoji": "🟡"},
            {"word": "green",    "zh": "綠色",  "zhuyin": "ㄌㄩˋ ㄙㄜˋ",     "emoji": "🟢"},
            {"word": "purple",   "zh": "紫色",  "zhuyin": "ㄗˇ ㄙㄜˋ",       "emoji": "🟣"},
            {"word": "pink",     "zh": "粉紅色", "zhuyin": "ㄈㄣˇ ㄏㄨㄥˊ",   "emoji": "🩷"},
            {"word": "shirt",    "zh": "上衣",  "zhuyin": "ㄕㄤˋ ㄧ",        "emoji": "👕"},
            {"word": "pants",    "zh": "褲子",  "zhuyin": "ㄎㄨˋ ˙ㄗ",       "emoji": "👖"},
            {"word": "dress",    "zh": "洋裝",  "zhuyin": "ㄧㄤˊ ㄓㄨㄤ",    "emoji": "👗"},
            {"word": "shoes",    "zh": "鞋子",  "zhuyin": "ㄒㄧㄝˊ ˙ㄗ",     "emoji": "👟"},
            {"word": "hat",      "zh": "帽子",  "zhuyin": "ㄇㄠˋ ˙ㄗ",       "emoji": "🎩"},
            {"word": "socks",    "zh": "襪子",  "zhuyin": "ㄨㄚˋ ˙ㄗ",       "emoji": "🧦"},
            {"word": "jacket",   "zh": "夾克",  "zhuyin": "ㄐㄧㄚ ㄎㄜˋ",    "emoji": "🧥"},
            {"word": "scarf",    "zh": "圍巾",  "zhuyin": "ㄨㄟˊ ㄐㄧㄣ",    "emoji": "🧣"},
            {"word": "gloves",   "zh": "手套",  "zhuyin": "ㄕㄡˇ ㄊㄠˋ",     "emoji": "🧤"},
            {"word": "boots",    "zh": "靴子",  "zhuyin": "ㄒㄩㄝ ˙ㄗ",      "emoji": "👢"},
        ],
        "color": RGBColor(0xE9, 0x5E, 0xA0),
    },
    {
        "id": 7, "title": "What Time Is It?", "emoji": "⏰",
        "vocab": [
            {"word": "one",       "zh": "一",    "zhuyin": "ㄧ",              "emoji": "1️⃣"},
            {"word": "two",       "zh": "二",    "zhuyin": "ㄦˋ",             "emoji": "2️⃣"},
            {"word": "three",     "zh": "三",    "zhuyin": "ㄙㄢ",            "emoji": "3️⃣"},
            {"word": "four",      "zh": "四",    "zhuyin": "ㄙˋ",             "emoji": "4️⃣"},
            {"word": "five",      "zh": "五",    "zhuyin": "ㄨˇ",             "emoji": "5️⃣"},
            {"word": "six",       "zh": "六",    "zhuyin": "ㄌㄧㄡˋ",         "emoji": "6️⃣"},
            {"word": "seven",     "zh": "七",    "zhuyin": "ㄑㄧ",            "emoji": "7️⃣"},
            {"word": "eight",     "zh": "八",    "zhuyin": "ㄅㄚ",            "emoji": "8️⃣"},
            {"word": "nine",      "zh": "九",    "zhuyin": "ㄐㄧㄡˇ",         "emoji": "9️⃣"},
            {"word": "ten",       "zh": "十",    "zhuyin": "ㄕˊ",             "emoji": "🔟"},
            {"word": "clock",     "zh": "時鐘",  "zhuyin": "ㄕˊ ㄓㄨㄥ",     "emoji": "🕐"},
            {"word": "time",      "zh": "時間",  "zhuyin": "ㄕˊ ㄐㄧㄢ",     "emoji": "⌚"},
            {"word": "morning",   "zh": "早上",  "zhuyin": "ㄗㄠˇ ˙ㄕㄤ",    "emoji": "🌅"},
            {"word": "afternoon", "zh": "下午",  "zhuyin": "ㄒㄧㄚˋ ˙ㄨ",    "emoji": "☀️"},
            {"word": "night",     "zh": "晚上",  "zhuyin": "ㄨㄢˇ ˙ㄕㄤ",    "emoji": "🌙"},
            {"word": "o'clock",   "zh": "...點鐘","zhuyin": "˙ㄉㄧㄢ",        "emoji": "⏰"},
        ],
        "color": RGBColor(0x16, 0xA0, 0x85),
    },
    {
        "id": 8, "title": "What's the Weather?", "emoji": "🌤️",
        "vocab": [
            {"word": "sunny",    "zh": "晴天",  "zhuyin": "ㄑㄧㄥˊ ㄊㄧㄢ",   "emoji": "☀️"},
            {"word": "rainy",    "zh": "下雨",  "zhuyin": "ㄒㄧㄚˋ ㄩˇ",      "emoji": "🌧️"},
            {"word": "cloudy",   "zh": "陰天",  "zhuyin": "ㄧㄣ ㄊㄧㄢ",      "emoji": "☁️"},
            {"word": "windy",    "zh": "刮風",  "zhuyin": "ㄍㄨㄚ ㄈㄥ",      "emoji": "💨"},
            {"word": "snowy",    "zh": "下雪",  "zhuyin": "ㄒㄧㄚˋ ㄒㄩㄝˇ",  "emoji": "❄️"},
            {"word": "foggy",    "zh": "有霧",  "zhuyin": "ㄧㄡˇ ㄨˋ",        "emoji": "🌫️"},
            {"word": "cold",     "zh": "冷",    "zhuyin": "ㄌㄥˇ",            "emoji": "🥶"},
            {"word": "hot",      "zh": "熱",    "zhuyin": "ㄖㄜˋ",            "emoji": "🥵"},
            {"word": "warm",     "zh": "溫暖",  "zhuyin": "ㄨㄣ ㄋㄨㄢˇ",     "emoji": "🌤️"},
            {"word": "cool",     "zh": "涼爽",  "zhuyin": "ㄌㄧㄤˊ ㄕㄨㄤˇ",  "emoji": "🍃"},
            {"word": "umbrella", "zh": "雨傘",  "zhuyin": "ㄩˇ ㄙㄢˇ",        "emoji": "☂️"},
            {"word": "coat",     "zh": "外套",  "zhuyin": "ㄨㄞˋ ㄊㄠˋ",      "emoji": "🧥"},
            {"word": "rainbow",  "zh": "彩虹",  "zhuyin": "ㄘㄞˇ ㄏㄨㄥˊ",    "emoji": "🌈"},
            {"word": "storm",    "zh": "暴風雨", "zhuyin": "ㄅㄠˋ ㄈㄥ ㄩˇ",   "emoji": "⛈️"},
            {"word": "thunder",  "zh": "打雷",  "zhuyin": "ㄉㄚˇ ㄌㄟˊ",      "emoji": "⚡"},
            {"word": "weather",  "zh": "天氣",  "zhuyin": "ㄊㄧㄢ ㄑㄧˋ",     "emoji": "🌡️"},
        ],
        "color": RGBColor(0x27, 0x6F, 0xBF),
    },
]

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK  = RGBColor(0x2C, 0x3E, 0x50)
LIGHT_BG = RGBColor(0xF8, 0xF9, 0xFA)

W = Inches(13.33)
H = Inches(7.5)


def add_text(tf, text, size, bold=False, color=WHITE, align=PP_ALIGN.CENTER):
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_box(slide, left, top, width, height, text, size, bold=False,
            txt_color=WHITE, bg=None, align=PP_ALIGN.CENTER):
    box = slide.shapes.add_textbox(left, top, width, height)
    if bg:
        box.fill.solid()
        box.fill.fore_color.rgb = bg
    tf = box.text_frame
    tf.auto_size = None
    add_text(tf, text, size, bold=bold, color=txt_color, align=align)
    return box


def make_title_slide(prs):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_bg(slide, RGBColor(0x2C, 0x3E, 0x50))

    add_box(slide, Inches(1), Inches(1.8), Inches(11.33), Inches(1.2),
            "Grade 4 English Tutor", 52, bold=True)
    add_box(slide, Inches(1), Inches(3.1), Inches(11.33), Inches(0.8),
            "四年級英語教材  ·  Units 1–8  ·  單字總覽", 26)
    add_box(slide, Inches(1), Inches(4.2), Inches(11.33), Inches(0.7),
            "每單元 16 個單字　共 128 個核心詞彙", 20)


def make_unit_cover(prs, unit):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_bg(slide, unit["color"])

    add_box(slide, Inches(1), Inches(1.5), Inches(11.33), Inches(1.4),
            f"Unit {unit['id']}", 44, bold=True)
    add_box(slide, Inches(1), Inches(2.9), Inches(11.33), Inches(1.6),
            unit["title"], 56, bold=True)
    add_box(slide, Inches(1), Inches(4.7), Inches(11.33), Inches(0.8),
            f"Vocabulary  ·  {len(unit['vocab'])} words", 24)


def make_word_slide(prs, unit, vocab, idx, total):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_bg(slide, LIGHT_BG)

    # top color bar
    bar = slide.shapes.add_shape(1, 0, 0, W, Inches(0.55))
    bar.fill.solid()
    bar.fill.fore_color.rgb = unit["color"]
    bar.line.fill.background()
    tf = bar.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    run = tf.paragraphs[0].add_run()
    run.text = f"  Unit {unit['id']} · {unit['title']}   {idx}/{total}"
    run.font.size = Pt(16)
    run.font.color.rgb = WHITE

    # emoji
    add_box(slide, Inches(0.5), Inches(0.7), Inches(12.33), Inches(2.2),
            vocab["emoji"], 96, txt_color=DARK)

    # english word
    add_box(slide, Inches(0.5), Inches(2.85), Inches(12.33), Inches(1.5),
            vocab["word"], 72, bold=True, txt_color=unit["color"])

    # divider line via shape
    line = slide.shapes.add_shape(1, Inches(4.5), Inches(4.45), Inches(4.33), Emu(4))
    line.fill.solid()
    line.fill.fore_color.rgb = unit["color"]
    line.line.fill.background()

    # chinese
    add_box(slide, Inches(0.5), Inches(4.55), Inches(6.3), Inches(1.0),
            vocab["zh"], 40, bold=True, txt_color=DARK)

    # zhuyin
    add_box(slide, Inches(6.8), Inches(4.55), Inches(6.0), Inches(1.0),
            vocab["zhuyin"], 32, txt_color=RGBColor(0x7F, 0x8C, 0x8D))

    # bottom word number pill
    pill = slide.shapes.add_shape(9, Inches(6.0), Inches(6.9), Inches(1.33), Inches(0.45))
    pill.fill.solid()
    pill.fill.fore_color.rgb = unit["color"]
    pill.line.fill.background()
    tf2 = pill.text_frame
    tf2.paragraphs[0].alignment = PP_ALIGN.CENTER
    r2 = tf2.paragraphs[0].add_run()
    r2.text = f"{idx} / {total}"
    r2.font.size = Pt(16)
    r2.font.color.rgb = WHITE


def make_unit_summary(prs, unit):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_bg(slide, LIGHT_BG)

    bar = slide.shapes.add_shape(1, 0, 0, W, Inches(0.55))
    bar.fill.solid()
    bar.fill.fore_color.rgb = unit["color"]
    bar.line.fill.background()
    tf = bar.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    run = tf.paragraphs[0].add_run()
    run.text = f"  Unit {unit['id']} · {unit['title']}   Summary"
    run.font.size = Pt(16)
    run.font.color.rgb = WHITE

    add_box(slide, Inches(0.4), Inches(0.65), Inches(12.5), Inches(0.7),
            f"Unit {unit['id']} 單字總整理", 28, bold=True, txt_color=unit["color"])

    cols = 4
    rows = 4
    cell_w = Inches(3.1)
    cell_h = Inches(1.4)
    start_x = Inches(0.2)
    start_y = Inches(1.4)

    for i, v in enumerate(unit["vocab"]):
        col = i % cols
        row = i // cols
        x = start_x + col * cell_w
        y = start_y + row * cell_h

        cell = slide.shapes.add_shape(9, x, y, cell_w - Inches(0.1), cell_h - Inches(0.1))
        cell.fill.solid()
        cell.fill.fore_color.rgb = WHITE
        cell.line.color.rgb = unit["color"]
        cell.line.width = Pt(1.5)

        tf3 = cell.text_frame
        tf3.word_wrap = True

        p1 = tf3.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = f"{v['emoji']}  {v['word']}"
        r1.font.size = Pt(18)
        r1.font.bold = True
        r1.font.color.rgb = DARK

        p2 = tf3.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = v["zh"]
        r2.font.size = Pt(14)
        r2.font.color.rgb = unit["color"]


def main():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    make_title_slide(prs)

    for unit in UNITS:
        make_unit_cover(prs, unit)
        total = len(unit["vocab"])
        for idx, v in enumerate(unit["vocab"], 1):
            make_word_slide(prs, unit, v, idx, total)
        make_unit_summary(prs, unit)

    out = r"C:\Users\yuchi\grade4-english-tutor\Grade4_English_Vocab.pptx"
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    main()
