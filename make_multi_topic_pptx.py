# -*- coding: utf-8 -*-
"""Multi-Topic Review — 16-頁跨主題複習簡報產生器。
主題：天氣 / 時間 / 情緒 / 教室用品 / 介系詞 in·on·at
主色：黃色系（amber）；文字 navy #2C3E50；無粉紅。
輸出：Multi_Topic_Review.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------- 主題色 ----------
THEME   = RGBColor(0xFF, 0xC1, 0x07)   # amber 500
THEME_D = RGBColor(0xB0, 0x73, 0x09)   # 深琥珀
LIGHT   = RGBColor(0xFF, 0xFA, 0xE5)   # 米色背景
CARD_BD = RGBColor(0xF1, 0xDA, 0x9B)   # 暖棕邊框
DARK    = RGBColor(0x2C, 0x3E, 0x50)   # navy slate
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GRAY    = RGBColor(0x7F, 0x8C, 0x8D)

W = Inches(13.333)
H = Inches(7.5)

OUT = r"C:\Users\yuchi\grade4-english-tutor\Multi_Topic_Review.pptx"


# ---------- 共用工具 ----------
def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def textbox(slide, x, y, w, h, lines, anchor=MSO_ANCHOR.MIDDLE, wrap=True):
    """lines: list of (text, size, bold, color, align)"""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for i, (text, size, bold, color, align) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return box


def rrect(slide, x, y, w, h, fill, border=None, bw=2.0):
    sh = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    sh.shadow.inherit = False
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if border is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = border
        sh.line.width = Pt(bw)
    return sh


def rect(slide, x, y, w, h, fill, border=None, bw=2.0):
    sh = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    sh.shadow.inherit = False
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if border is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = border
        sh.line.width = Pt(bw)
    return sh


def topbar(slide, label):
    bar = rect(slide, 0, 0, 13.333, 0.68, THEME)
    tf = bar.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.size = Pt(32)
    r.font.bold = True
    r.font.color.rgb = DARK


def word_card(slide, x, y, w, h, emoji, word, chinese, large=False):
    """白底字卡：emoji + 英文 + 中文"""
    em_pt = 50 if large else 36
    en_pt = 30 if large else 22
    zh_pt = 22 if large else 16

    card = rrect(slide, x, y, w, h, WHITE, CARD_BD, bw=2.5)
    tf = card.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r = p1.add_run()
    r.text = emoji
    r.font.size = Pt(em_pt)

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r = p2.add_run()
    r.text = word
    r.font.size = Pt(en_pt)
    r.font.bold = True
    r.font.color.rgb = THEME_D

    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    r = p3.add_run()
    r.text = chinese
    r.font.size = Pt(zh_pt)
    r.font.color.rgb = DARK


def grid_6(slide, words):
    """2 欄 × 3 列，共 6 張字卡"""
    mx, gx = 0.3, 0.26
    cw = (13.333 - mx * 2 - gx) / 2   # ≈ 6.237"
    ty, gy = 0.78, 0.16
    ch = (7.5 - ty - 0.12 - gy * 2) / 3   # ≈ 2.093"
    for i, (emoji, word, zh) in enumerate(words):
        col = i % 2
        row = i // 2
        x = mx + col * (cw + gx)
        y = ty + row * (ch + gy)
        word_card(slide, x, y, cw, ch, emoji, word, zh, large=True)


def grid_8(slide, words):
    """2 欄 × 4 列，共 8 張字卡"""
    mx, gx = 0.3, 0.26
    cw = (13.333 - mx * 2 - gx) / 2   # ≈ 6.237"
    ty, gy = 0.78, 0.13
    ch = (7.5 - ty - 0.10 - gy * 3) / 4   # ≈ 1.533"
    for i, (emoji, word, zh) in enumerate(words):
        col = i % 2
        row = i // 2
        x = mx + col * (cw + gx)
        y = ty + row * (ch + gy)
        word_card(slide, x, y, cw, ch, emoji, word, zh, large=False)


def pattern_box(slide, q_line, a_line, zh_line, example_line=""):
    """句型大框（米色底 + 深琥珀框線）"""
    box = rrect(slide, 0.45, 0.98, 12.433, 4.6, LIGHT, THEME_D, bw=3.5)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r = p1.add_run()
    r.text = q_line
    r.font.size = Pt(42)
    r.font.bold = True
    r.font.color.rgb = DARK

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r = p2.add_run()
    r.text = a_line
    r.font.size = Pt(42)
    r.font.bold = True
    r.font.color.rgb = THEME_D

    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    r = p3.add_run()
    r.text = zh_line
    r.font.size = Pt(24)
    r.font.bold = False
    r.font.color.rgb = GRAY

    if example_line:
        textbox(slide, 0.45, 5.72, 12.433, 0.9,
                [(example_line, 20, False, THEME_D, PP_ALIGN.CENTER)])


# ---------- 16 張投影片 ----------

def s01_cover(prs):
    s = blank(prs)
    set_bg(s, THEME)

    # 裝飾大圓
    for ox, oy, od in [(10.5, -1.2, 4.2), (-1.8, 4.8, 3.8)]:
        c = s.shapes.add_shape(MSO_SHAPE.OVAL,
                               Inches(ox), Inches(oy), Inches(od), Inches(od))
        c.shadow.inherit = False
        c.fill.solid(); c.fill.fore_color.rgb = THEME_D
        c.line.fill.background()

    # 頂部 emoji 裝飾
    textbox(s, 0, 0.55, 13.333, 0.9,
            [("⛅   🕐   😊   🏫   📍", 48, False, DARK, PP_ALIGN.CENTER)])

    # 主標題框
    title_box = rrect(s, 1.0, 1.6, 11.333, 2.35, THEME_D, None)
    tf = title_box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Multi-Topic Review"
    r.font.size = Pt(54)
    r.font.bold = True
    r.font.color.rgb = WHITE

    # 副標題
    textbox(s, 1.0, 4.1, 11.333, 0.8,
            [("Grade 4 English · 四年級英語複習", 28, False, DARK, PP_ALIGN.CENTER)])

    # 主題標籤列
    tag = rrect(s, 1.0, 5.05, 11.333, 0.78, WHITE, CARD_BD, bw=2.0)
    tf2 = tag.text_frame
    tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "天氣  ·  時間  ·  情緒  ·  教室用品  ·  介系詞  in / on / at"
    r2.font.size = Pt(22)
    r2.font.color.rgb = THEME_D

    # 網址
    textbox(s, 0, 6.15, 13.333, 0.65,
            [("https://Yuchiaoniu.github.io/english/", 17, False, DARK, PP_ALIGN.CENTER)])


def s02_weather_vocab(prs):
    s = blank(prs)
    set_bg(s, LIGHT)
    topbar(s, "⛅  Weather  天氣單字")
    grid_6(s, [
        ("☀️",  "sunny",  "晴天"),
        ("⛅",  "cloudy", "多雲"),
        ("🌧️", "rainy",  "下雨"),
        ("💨",  "windy",  "颳風"),
        ("❄️",  "snowy",  "下雪"),
        ("🥶",  "cold",   "寒冷"),
    ])


def s03_weather_pattern(prs):
    s = blank(prs)
    set_bg(s, LIGHT)
    topbar(s, "⛅  Weather  天氣句型")
    pattern_box(
        s,
        "What's the weather like?",
        "It's  ___.",
        "（天氣怎麼樣？／ 今天天氣是 ___ 。）",
        "☀️ It's sunny.    ⛅ It's cloudy.    🌧️ It's rainy.    💨 It's windy."
    )


def s04_time_vocab(prs):
    s = blank(prs)
    set_bg(s, LIGHT)
    topbar(s, "🕐  Time  時間單字")
    grid_6(s, [
        ("🌅",  "morning",   "早上"),
        ("☀️",  "afternoon", "下午"),
        ("🌆",  "evening",   "傍晚"),
        ("🌙",  "night",     "晚上"),
        ("🕐",  "o'clock",   "點鐘"),
        ("🕧",  "half past", "半（30 分）"),
    ])


def s05_time_pattern(prs):
    s = blank(prs)
    set_bg(s, LIGHT)
    topbar(s, "🕐  Time  時間句型")
    pattern_box(
        s,
        "What time is it?",
        "It's  ___  o'clock.",
        "（現在幾點了？／ 現在是 ___ 點。）",
        "🕒 It's three o'clock.    🕠 It's half past five."
    )


def s06_feelings_vocab(prs):
    s = blank(prs)
    set_bg(s, LIGHT)
    topbar(s, "😊  Feelings  情緒感覺")
    grid_8(s, [
        ("😄", "happy",     "開心"),
        ("😢", "sad",       "難過"),
        ("😠", "angry",     "生氣"),
        ("😨", "scared",    "害怕"),
        ("🤩", "excited",   "興奮"),
        ("😴", "tired",     "疲倦"),
        ("🏆", "proud",     "驕傲"),
        ("😲", "surprised", "驚訝"),
    ])


def s07_feelings_pattern(prs):
    s = blank(prs)
    set_bg(s, LIGHT)
    topbar(s, "😊  Feelings  情緒句型")
    pattern_box(
        s,
        "How do you feel?",
        "I feel  ___.",
        "（你感覺怎麼樣？／ 我感覺 ___ 。）",
        "😄 I feel happy.    😢 I feel sad.    😠 I feel angry."
    )


def s08_classroom_vocab(prs):
    s = blank(prs)
    set_bg(s, LIGHT)
    topbar(s, "🏫  Classroom  教室用品")
    grid_8(s, [
        ("🪑", "chair",  "椅子"),
        ("📋", "board",  "黑板"),
        ("🪟", "window", "窗戶"),
        ("🚪", "door",   "門"),
        ("📚", "book",   "書"),
        ("✏️", "pen",    "筆"),
        ("📏", "ruler",  "尺"),
        ("🎒", "bag",    "書包"),
    ])


def s09_classroom_pattern(prs):
    s = blank(prs)
    set_bg(s, LIGHT)
    topbar(s, "🏫  Classroom  教室句型")
    pattern_box(
        s,
        "Where is the  ___?",
        "It is  ___  the  ___.",
        "（___ 在哪裡？／ 它在 ___ 的 ___ 。）",
        "📚 The book is on the desk.    ✏️ The pen is in the bag."
    )


def s10_prep_intro(prs):
    s = blank(prs)
    set_bg(s, LIGHT)
    topbar(s, "📍  Prepositions  介系詞 in / on / at")

    cols = [
        ("AT", "📍", "點（Point）",
         "一個明確的地點",
         "at home", "at school", "at the bus stop"),
        ("ON", "🛣️", "線／面（Line / Surface）",
         "在某個表面或路線上",
         "on the bus", "on the road", "on the desk"),
        ("IN", "📦", "面／空間（Area / Space）",
         "在某個空間或範圍裡面",
         "in the box", "in the classroom", "in Taiwan"),
    ]

    cw = 4.044
    mx = 0.35
    gx = 0.25
    cy = 0.82
    ch = 6.38

    for i, (prep, sym, zh_title, zh_desc, ex1, ex2, ex3) in enumerate(cols):
        x = mx + i * (cw + gx)
        card = rrect(s, x, cy, cw, ch, WHITE, CARD_BD, bw=2.5)
        tf = card.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        rows = [
            (prep,     48, True,  THEME_D, PP_ALIGN.CENTER),
            (sym,      40, False, DARK,    PP_ALIGN.CENTER),
            (zh_title, 22, True,  DARK,    PP_ALIGN.CENTER),
            (zh_desc,  16, False, GRAY,    PP_ALIGN.CENTER),
            (ex1,      20, False, THEME_D, PP_ALIGN.CENTER),
            (ex2,      20, False, THEME_D, PP_ALIGN.CENTER),
            (ex3,      20, False, THEME_D, PP_ALIGN.CENTER),
        ]
        for j, (text, size, bold, color, align) in enumerate(rows):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.alignment = align
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color


def example_slide(prs, bar_label, cards):
    """AT / ON / IN 例句頁，3 張大例句卡"""
    s = blank(prs)
    set_bg(s, LIGHT)
    topbar(s, bar_label)

    cw = 4.044
    mx = 0.35
    gx = 0.25
    cy = 0.82
    ch = 6.38

    for i, (emoji, en, zh) in enumerate(cards):
        x = mx + i * (cw + gx)
        card = rrect(s, x, cy, cw, ch, WHITE, CARD_BD, bw=2.5)
        tf = card.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        data = [
            (emoji, 64, False, DARK,    PP_ALIGN.CENTER),
            (en,    30, True,  THEME_D, PP_ALIGN.CENTER),
            (zh,    22, False, DARK,    PP_ALIGN.CENTER),
        ]
        for j, (text, size, bold, color, align) in enumerate(data):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.alignment = align
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color


def s11_at_examples(prs):
    example_slide(prs, "📍  AT  點狀位置（明確的一個地點）", [
        ("🏠", "at home",        "在家"),
        ("🏫", "at school",      "在學校"),
        ("🚌", "at the bus stop","在公車站"),
    ])


def s12_on_examples(prs):
    example_slide(prs, "🛣️  ON  表面 / 路線上", [
        ("🛣️", "on the road",  "在路上"),
        ("🚌", "on the bus",   "在公車上"),
        ("📚", "on the shelf", "在書架上"),
    ])


def s13_in_examples(prs):
    example_slide(prs, "📦  IN  空間 / 範圍裡", [
        ("🏫", "in the classroom", "在教室裡"),
        ("📦", "in the box",       "在箱子裡"),
        ("🌏", "in Taiwan",        "在台灣"),
    ])


def s14_fillin(prs):
    s = blank(prs)
    set_bg(s, LIGHT)
    topbar(s, "✏️  Fill in the Blank  填空練習（in / on / at）")

    qs = [
        ("1", "The cat is ___  the box.",
         "貓在箱子裡。", "in"),
        ("2", "She is ___  school.",
         "她在學校。", "at"),
        ("3", "The book is ___  the desk.",
         "書在桌子上。", "on"),
        ("4", "We are ___  Taiwan.",
         "我們在台灣。", "in"),
    ]

    qh = 1.44
    qx = 0.42
    qw = 12.493
    ty = 0.82
    gy = 0.12

    for i, (num, q, zh, ans) in enumerate(qs):
        y = ty + i * (qh + gy)
        card = rrect(s, qx, y, qw, qh, WHITE, CARD_BD, bw=1.8)
        tf = card.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        rows = [
            (f"{num}.  {q}   →  Answer :  ( {ans} )", 28, True,  DARK,    PP_ALIGN.LEFT),
            (f"       {zh}",                            20, False, GRAY,    PP_ALIGN.LEFT),
        ]
        for j, (text, size, bold, color, align) in enumerate(rows):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.alignment = align
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color


def s15_quiz(prs):
    s = blank(prs)
    set_bg(s, LIGHT)
    topbar(s, "🌟  Review Quiz  複習測驗")

    qs = [
        ("1", "What's the weather like?   ___ cloudy.",
         "(A) It's   (B) I'm   (C) They're",  "A"),
        ("2", "How do you feel?   I feel ___.",
         "(A) happy   (B) morning   (C) school", "A"),
        ("3", "The pen is ___ the desk.",
         "(A) in   (B) on   (C) at",             "B"),
        ("4", "___ is it?   It's three o'clock.",
         "(A) Where   (B) What time   (C) How",  "B"),
    ]

    qh = 1.38
    qx = 0.42
    qw = 12.493
    ty = 0.82
    gy = 0.12

    for i, (num, q, opts, ans) in enumerate(qs):
        y = ty + i * (qh + gy)
        card = rrect(s, qx, y, qw, qh, WHITE, CARD_BD, bw=1.8)
        tf = card.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        rows = [
            (f"{num}.  {q}",          26, True,  DARK,    PP_ALIGN.LEFT),
            (f"   {opts}   ✓ {ans}", 20, False, THEME_D, PP_ALIGN.LEFT),
        ]
        for j, (text, size, bold, color, align) in enumerate(rows):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.alignment = align
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color


def s16_end(prs):
    s = blank(prs)
    set_bg(s, THEME)

    # 裝飾圓
    for ox, oy, od in [(10.5, -1.2, 4.2), (-1.8, 4.8, 3.8)]:
        c = s.shapes.add_shape(MSO_SHAPE.OVAL,
                               Inches(ox), Inches(oy), Inches(od), Inches(od))
        c.shadow.inherit = False
        c.fill.solid(); c.fill.fore_color.rgb = THEME_D
        c.line.fill.background()

    textbox(s, 0, 1.0, 13.333, 1.1,
            [("🌟   🌟   🌟", 56, False, DARK, PP_ALIGN.CENTER)])

    title_box = rrect(s, 1.4, 2.25, 10.533, 2.1, THEME_D, None)
    tf = title_box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Great Job!  你做到了！"
    r.font.size = Pt(50)
    r.font.bold = True
    r.font.color.rgb = WHITE

    textbox(s, 0, 4.55, 13.333, 0.75,
            [("你已完成所有主題的複習！Keep it up!", 26, False, DARK, PP_ALIGN.CENTER)])

    textbox(s, 0, 5.4, 13.333, 0.65,
            [("https://Yuchiaoniu.github.io/english/", 18, False, DARK, PP_ALIGN.CENTER)])


# ---------- 主程式 ----------
def make():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    s01_cover(prs)
    s02_weather_vocab(prs)
    s03_weather_pattern(prs)
    s04_time_vocab(prs)
    s05_time_pattern(prs)
    s06_feelings_vocab(prs)
    s07_feelings_pattern(prs)
    s08_classroom_vocab(prs)
    s09_classroom_pattern(prs)
    s10_prep_intro(prs)
    s11_at_examples(prs)
    s12_on_examples(prs)
    s13_in_examples(prs)
    s14_fillin(prs)
    s15_quiz(prs)
    s16_end(prs)

    prs.save(OUT)
    print(f"Saved  : {OUT}")
    print(f"Slides : {len(prs.slides)}")


if __name__ == "__main__":
    make()
